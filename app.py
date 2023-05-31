import boto3
import json
from pptx import Presentation
import shutil
import rt
import os
import datetime

def search_and_replace(search_str, repl_str, file):
    """"search and replace text in PowerPoint while preserving formatting"""
    loop = len(search_str)
    prs = Presentation(file)
    for i in range(loop):
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if(shape.text.find(search_str[i]))!=-1:
                        text_frame = shape.text_frame
                        cur_text = text_frame.paragraphs[0].runs[0].text
                        new_text = cur_text.replace(str(search_str[i]), str(repl_str[i]))
                        text_frame.paragraphs[0].runs[0].text = new_text
    prs.save(file)

def handler(event, context):
    customer_map = {
        "ADOT": "Arizona Department of Transportation",
        "ASRS": "Arizona State Retirement System",
        "SFO": "San Francisco International Airport",
        "SDS": "San Domenico School",
        "SNHD": "Southern NV Health District",
        "DBS":  "DB Schenker",
        "FMS":  "Fenix Marine Services",
        "JWA":  "John Wayne Airport",
        "M451": "Mosaic451",
        "TJPA": "TJPA"
    }
    customer_str = ["first_day","last_day","mth_year","customer_acr","customer_name","total_logs_ingested","security_alerts","cases_opened","cases_escalated"]

    cross_account_external_id = {
        "DBS" : f"{os.environ['dbs_role_external_id']}"
    }

    customer_accounts = {
        "M451" : "",
        "DBS"  : f"{os.environ['dbs_role']}"
    }


    # SOCRT Login
    tracker = rt.Rt('https://socrt.mosaic451.com/REST/1.0/', os.environ['rtuser'], os.environ['rtpasswd'])
    tracker.login()

# Search through the s3 buckets for the .json data

    for accounts in customer_accounts.keys():
        # Assume role for other customers AWS accounts
        if customer_accounts[accounts] != "":
            sts_client = boto3.client('sts')

            assumed_role_object=sts_client.assume_role(
                RoleArn=customer_accounts[accounts],
                RoleSessionName="AssumeRoleSession1",
                ExternalId=cross_account_external_id[accounts]
            )
            credentials=assumed_role_object['Credentials']

            s3 = boto3.resource(
                's3',
                aws_access_key_id=credentials['AccessKeyId'],
                aws_secret_access_key=credentials['SecretAccessKey'],
                aws_session_token=credentials['SessionToken'],
            )
        # If CrossAccount is blank, use default creds.
        else:
            s3 = boto3.resource('s3')

        for bucket in s3.buckets.all():
            if bucket.name.endswith("ngs-report-logs"):
                for obj in bucket.objects.filter(Prefix=f"compiled_reports/year={event['year']}/month={event['month']}/dashboard_report.json"):
                    print('{0}:{1}'.format(bucket.name, obj.key))
                    content_object = s3.Object(bucket.name, obj.key)
                    file_content = content_object.get()['Body'].read().decode('utf-8')
                    report_data = json.loads(file_content)               

                    rolled_up_log_count = {}
                    rolled_up_alert_count = {}
                    rolled_up_case_count = {}

                    for log_csv in report_data['log_csv']:        
                        if log_csv in rolled_up_log_count.keys():            
                            rolled_up_log_count[log_csv] = rolled_up_log_count[log_csv] + int(report_data['log_csv'][log_csv]['count'].replace('.0',''))        
                        else:            
                            rolled_up_log_count[log_csv] = int(report_data['log_csv'][log_csv]['count'].replace('.0',''))

                    for alert_csv in report_data['alert_csv']:        
                        if alert_csv == 'Alert Category' or alert_csv == 'Silent Log - Dbs':            
                            continue        
                        if alert_csv in rolled_up_alert_count.keys():            
                            rolled_up_alert_count[alert_csv] = rolled_up_alert_count[alert_csv] + int(report_data['alert_csv'][alert_csv]['count_of_alerts'].replace('.0',''))        
                        else:            
                            rolled_up_alert_count[alert_csv] = int(report_data['alert_csv'][alert_csv]['count_of_alerts'].replace('.0',''))

                    for case_csv in report_data['case_csv']:      
                        if case_csv in rolled_up_case_count.keys():            
                            rolled_up_case_count[case_csv] = rolled_up_case_count[case_csv] + int(report_data['case_csv'][case_csv])       
                        else:            
                            rolled_up_case_count[case_csv] = int(report_data['case_csv'][case_csv])

                    customer_acr=bucket.name.split("-")[0].upper()

                    year=int(event['year'])
                    month=int(event['month'])
                    day=int(event['day'])

                    query_first_day=datetime.date(year,month,1)
                    query_last_day=datetime.date(year,month,day)
                    first_day=query_first_day.strftime("%B %d, %Y")
                    last_day=query_last_day.strftime("%B %d, %Y")
                    mth_year=query_first_day.strftime("%B %Y")

                    rolled_up_cases_escalated=len(list(map(lambda x: x['id'], tracker.search(Queue=f"{customer_acr}-Investigations",raw_query=f"Created > '{query_first_day}' AND Created < '{query_last_day}'"))))
                    customer_info = []
                    customer_info.append(first_day)
                    customer_info.append(last_day)
                    customer_info.append(mth_year)
                    customer_info.append(customer_acr)
                    customer_info.append(customer_map[customer_acr])
                    customer_info.append(f'{sum(rolled_up_log_count.values()):,}')
                    customer_info.append(f'{sum(rolled_up_alert_count.values()):,}')
                    customer_info.append(f'{sum(rolled_up_case_count.values()):,}')
                    customer_info.append(f'{rolled_up_cases_escalated:,}')

                    report_name=f"{customer_acr}.pptx"
                    report_file_path=f"/tmp/{bucket.name}.pptx"
                    
                    shutil.copyfile('template.pptx', report_file_path)

                    search_and_replace(customer_str,customer_info,report_file_path)

                    # Upload file to S3 bucket
                    
                    s3_client = boto3.client('s3')
                    s3_client.upload_file(report_file_path,"m451-tenant-reports",f"{event['year']}/{event['month']}/{report_name}")

        # SOCRT LOGOUT
    tracker.logout()

    return True