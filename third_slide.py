import re
import os
import glob
import csv
from pptx.util import Pt
from pptx import Presentation


BASE_DIR = os.path.dirname(os.path.abspath(__file__))


def formatINR(number):
    s, *d = str(number).partition(".")
    r = ",".join([s[x-2:x] for x in range(-3, -len(s), -2)][::-1] + [s[-3:]])
    return "".join([r] + d)


def search_or_replace(repl_str, file, tmp_var):

    prs = Presentation(file)
    i = 0
    row_tmp = 0
    tmp_list = []
    tmp1_list = []
    slide = prs.slides[2]
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                if row_tmp == 0:
                    row_tmp += 1
                    pass
                else:
                    for cell in row.cells:
                        for kk in range(len(repl_str)):
                            if cell.text.replace(",", "").replace(" ", "").lower() == repl_str[kk]:
                                i = 0
                                for cell in row.cells:
                                    if i == 1:
                                        new_text = cell.text.replace(
                                            cell.text, str(formatINR(repl_str[kk + i])))
                                        cell.text = new_text
                                        i += 1
                                        for paragraph in cell.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                run.font.size = Pt(10)
                                    elif i == 2:
                                        new_text = cell.text.replace(
                                            cell.text, str(formatINR(repl_str[kk + i])+'%'))
                                        cell.text = new_text
                                        i += 1
                                        for paragraph in cell.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                run.font.size = Pt(10)
                                    else:
                                        new_text = cell.text.replace(
                                            cell.text, str(repl_str[kk + i]))
                                        cell.text = new_text
                                        i += 1
                                        for paragraph in cell.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                run.font.size = Pt(10)

                        tmp1_list.append(cell.text)

    ch_list = []
    for x in tmp1_list:
        if re.search('[a-zA-Z]+', x):
            ch_list.append(x)

    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                if row_tmp == 0:
                    row_tmp += 1
                    pass
                else:
                    for cell in row.cells:
                        if cell.text in ch_list:
                            i = 0
                            for cell in row.cells:
                                if cell.text in ch_list:
                                    pass
                                else:
                                    new_text = cell.text.replace(
                                        cell.text, str(0))
                                    cell.text = new_text
                                    i += 1
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.size = Pt(10)

    first_count = 0
    second_count = 0
    row_tmp = 0
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                if row_tmp == 0:
                    row_tmp += 1
                    pass
                else:
                    i = 0
                    for cell in row.cells:
                        if i == 0:
                            i += 1
                            pass
                        elif i == 1:
                            i += 1
                            first_count += float(cell.text.replace(",",
                                                 "").replace(" ", ""))
                        else:
                            i += 1
                            second_count += float(cell.text.replace(",",
                                                  "").replace(" ", "").replace("%", ""))

    i = 0
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        row_count = len(tbl.rows)
        col_count = len(tbl.columns)
        for c in range(0, col_count):
            cell = tbl.cell(row_count-1, c)
            if i == 0:
                i += 1
                pass
            elif i == 1:
                new_text = cell.text.replace(
                    cell.text, str(formatINR(first_count)))
                cell.text = new_text
                i += 1
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
            else:
                new_text = cell.text.replace(
                    cell.text, str(formatINR(second_count)+'%'))
                cell.text = new_text
                i += 1
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)

    out_file = "out-" + str(tmp_var) + ".pptx"
    prs.save(out_file)


def main():
    for filename in glob.glob(os.path.join(BASE_DIR, '*.pptx')):
        if filename.endswith('.pptx'):
            ppt_path = os.path.join(BASE_DIR, filename)

    tmp_var = 0
    for filename in glob.glob(os.path.join(BASE_DIR, '*.csv')):
        result = []
        tmp = {}
        if filename.endswith('.csv'):
            fullpath = os.path.join(BASE_DIR, filename)
            with open(fullpath, 'rt') as f:
                data = csv.reader(f)
                i = 0
                size_in_bytes = 0
                for row in data:
                    if i == 0:
                        pass
                    else:
                        size_in_bytes += float(row[5])
                        tmp[row[3]] = [size_in_bytes, size_in_bytes/100]

                    i += 1
                for key, value in tmp.items():
                    temp = [key, value[0], value[1]]
                    result.append(temp)

        flat_list = [item for sublist in result for item in sublist]

        search_or_replace(flat_list, ppt_path, tmp_var)
        tmp_var += 1


if __name__ == "__main__":
    main()